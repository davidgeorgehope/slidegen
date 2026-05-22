"""Render architecture-style slide specs as editable diagrams."""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from render import (
    DEFAULT_SOURCE_SIZE,
    SLIDE_H,
    SLIDE_W,
    add_contained_picture,
    add_picture_bbox,
    add_rect,
    add_text,
    emu,
    hex_to_rgb,
    no_line,
    px_box,
    px_to_emu_x,
    px_to_emu_y,
    set_solid_fill,
    spec_asset_path,
)


def rgb(theme: dict, key: str) -> str:
    return theme.get(key, key)


def text_box(slide, bbox, text, *, size, color, font, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    x, y, w, h = px_box(bbox)
    return add_text(
        slide,
        left=x,
        top=y,
        width=w,
        height=h,
        text=text,
        font_size=size,
        color_hex=color,
        bold=bold,
        align=align,
        anchor=anchor,
        font_name=font,
    )


def line(slide, x1, y1, x2, y2, color="285CDD", width=1.2, dashed=False, arrow=False):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        px_to_emu_x(x1),
        px_to_emu_y(y1),
        px_to_emu_x(x2),
        px_to_emu_y(y2),
    )
    shape.line.color.rgb = hex_to_rgb(color)
    shape.line.width = Pt(width)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        shape.line.end_arrowhead = True
    return shape


def circle_badge(slide, cx, cy, d, theme, kind):
    x = cx - d / 2
    y = cy - d / 2
    add_rect(
        slide,
        left=px_to_emu_x(x),
        top=px_to_emu_y(y),
        width=px_to_emu_x(d),
        height=px_to_emu_y(d),
        fill_hex="FFFFFF",
        line_hex="D5E3EF",
        line_pt=1.0,
        radius=0.5,
    )
    inner = add_rect(
        slide,
        left=px_to_emu_x(x + 9),
        top=px_to_emu_y(y + 9),
        width=px_to_emu_x(d - 18),
        height=px_to_emu_y(d - 18),
        fill_hex=theme["accent_color"],
        line_hex=None,
        radius=0.5,
    )
    if kind == "lock":
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px_to_emu_x(cx - 10), px_to_emu_y(cy - 2), px_to_emu_x(20), px_to_emu_y(17))
        set_solid_fill(sh, "FFFFFF")
        no_line(sh)
        arc = slide.shapes.add_shape(MSO_SHAPE.ARC, px_to_emu_x(cx - 9), px_to_emu_y(cy - 15), px_to_emu_x(18), px_to_emu_y(20))
        arc.line.color.rgb = hex_to_rgb("FFFFFF")
        arc.line.width = Pt(2)
    elif kind == "shield":
        sh = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, px_to_emu_x(cx - 10), px_to_emu_y(cy - 13), px_to_emu_x(20), px_to_emu_y(25))
        set_solid_fill(sh, "FFFFFF")
        no_line(sh)
    elif kind == "target":
        for r in (20, 12, 4):
            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, px_to_emu_x(cx - r / 2), px_to_emu_y(cy - r / 2), px_to_emu_x(r), px_to_emu_y(r))
            circ.fill.background()
            circ.line.color.rgb = hex_to_rgb("FFFFFF")
            circ.line.width = Pt(1.4)
    elif kind == "globe":
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, px_to_emu_x(cx - 14), px_to_emu_y(cy - 14), px_to_emu_x(28), px_to_emu_y(28))
        circ.fill.background()
        circ.line.color.rgb = hex_to_rgb("FFFFFF")
        circ.line.width = Pt(1.5)
        line(slide, cx - 13, cy, cx + 13, cy, "FFFFFF", 1.0)
        line(slide, cx, cy - 13, cx, cy + 13, "FFFFFF", 1.0)
    return inner


def device_icon(slide, cx, cy, kind, theme):
    if kind == "laptop":
        add_rect(slide, left=px_to_emu_x(cx - 22), top=px_to_emu_y(cy - 16), width=px_to_emu_x(44), height=px_to_emu_y(28), no_fill=True, line_hex=theme["accent_color"], line_pt=1.4, radius=0.04)
        line(slide, cx - 28, cy + 17, cx + 28, cy + 17, theme["accent_color"], 1.4)
    elif kind == "phone":
        add_rect(slide, left=px_to_emu_x(cx - 10), top=px_to_emu_y(cy - 22), width=px_to_emu_x(20), height=px_to_emu_y(44), no_fill=True, line_hex=theme["accent_color"], line_pt=1.4, radius=0.12)
    else:
        for dx in (-13, 13):
            slide.shapes.add_shape(MSO_SHAPE.OVAL, px_to_emu_x(cx + dx - 8), px_to_emu_y(cy - 19), px_to_emu_x(16), px_to_emu_y(16)).line.color.rgb = hex_to_rgb(theme["accent_color"])
            add_rect(slide, left=px_to_emu_x(cx + dx - 13), top=px_to_emu_y(cy + 1), width=px_to_emu_x(26), height=px_to_emu_y(22), no_fill=True, line_hex=theme["accent_color"], line_pt=1.4, radius=0.4)


def render_header(slide, spec, theme, fonts, base_dir):
    header = spec["header"]
    logo_path = base_dir / header.get("logo_image", "")
    if logo_path.exists() and logo_path.is_file():
        add_contained_picture(slide, logo_path, left=Inches(0.36), top=Inches(0.24), width=Inches(1.8), height=Inches(0.38))
    text_box(slide, [1380, 31, 94, 49], header["right_code"], size=12, color=theme["accent_color"], font=fonts["body"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, left=px_to_emu_x(1360), top=px_to_emu_y(29), width=px_to_emu_x(115), height=px_to_emu_y(50), no_fill=True, line_hex="8DB1F2", line_pt=1.0, radius=0.5)
    text_box(slide, [1500, 31, 145, 49], header["right_logo_placeholder"], size=12, color=theme["accent_color"], font=fonts["body"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, left=px_to_emu_x(1488), top=px_to_emu_y(29), width=px_to_emu_x(150), height=px_to_emu_y(50), no_fill=True, line_hex="8DB1F2", line_pt=1.0, radius=0.5)


def render_left_panel(slide, spec, theme, fonts):
    lp = spec["left_panel"]
    text_box(slide, [46, 108, 300, 22], lp["eyebrow"], size=11, color=theme["accent_color"], font=fonts["body"], bold=True)
    for i, line_text in enumerate(lp["title_lines"]):
        color = theme["accent_color"] if i == len(lp["title_lines"]) - 1 else theme["title_color"]
        text_box(slide, [46, 148 + i * 58, 410, 56], line_text, size=27, color=color, font=fonts["title"])
    text_box(slide, [46, 354, 368, 52], lp["body"], size=12, color=theme["body_color"], font=fonts["body"])
    line(slide, 46, 329, 111, 329, theme["accent_color"], 2)

    y_positions = [458, 550, 646, 754]
    for idx, (step, cy) in enumerate(zip(lp["steps"], y_positions)):
        circle_badge(slide, 76, cy, 58, theme, step["icon"])
        if idx < len(y_positions) - 1:
            line(slide, 76, cy + 29, 76, y_positions[idx + 1] - 29, theme["accent_color"], 1.0, dashed=True)
        text_box(slide, [124, cy - 27, 280, 26], step["name"], size=13, color=theme["accent_color"], font=fonts["title"], bold=True)
        text_box(slide, [124, cy, 295, 62], step["body"], size=10, color=theme["body_color"], font=fonts["body"])


def render_devices(slide, spec, theme, fonts):
    add_rect(slide, left=px_to_emu_x(577), top=px_to_emu_y(98), width=px_to_emu_x(774), height=px_to_emu_y(82), fill_hex="FFFFFF", line_hex=theme["card_border"], line_pt=1.0, radius=0.12)
    text_box(slide, [613, 129, 145, 24], "USERS & DEVICES", size=10, color=theme["accent_color"], font=fonts["body"], bold=True)
    xs = [842, 1037, 1225]
    for x, device in zip(xs, spec["devices"]):
        device_icon(slide, x, 139, device["icon"], theme)
        text_box(slide, [x + 40, 121, 95, 42], device["label"], size=8, color=theme["body_color"], font=fonts["body"], bold=True)
        line(slide, x + 40, 174, x + 40, 202, theme["accent_color"], 1.0, dashed=True)


def render_logo(slide, spec, base_dir, name, fallback_bbox, fallback_text, theme, fonts):
    path = spec_asset_path(spec, base_dir, name)
    if path:
        left, top, width, height = px_box(fallback_bbox)
        add_contained_picture(slide, path, left=left, top=top, width=width, height=height)
        return
    x, y, w, h = fallback_bbox
    add_rect(slide, left=px_to_emu_x(x), top=px_to_emu_y(y), width=px_to_emu_x(w), height=px_to_emu_y(h), fill_hex="FFFFFF", line_hex=theme["card_border"], radius=0.12)
    text_box(slide, fallback_bbox, fallback_text, size=9, color=theme["body_color"], font=fonts["body"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def bbox_center_px(bbox):
    x, y, w, h = bbox
    return x + w / 2, y + h / 2


def bbox_contains_point(bbox, point, pad=8):
    x, y, w, h = bbox
    px, py = point
    return x - pad <= px <= x + w + pad and y - pad <= py <= y + h + pad


def asset_or_fallback_bbox(spec, name, fallback_bbox, container_bbox=None):
    asset = spec.get("assets", {}).get(name, {})
    if isinstance(asset, dict) and asset.get("bbox"):
        if container_bbox and not bbox_contains_point(container_bbox, bbox_center_px(asset["bbox"])):
            return fallback_bbox
        return asset["bbox"]
    return fallback_bbox


def render_layers(slide, spec, theme, fonts, base_dir):
    labels = {item["name"]: item["match"] for item in spec.get("logo_assets", [])}

    for layer in spec["layers"]:
        x, y, w, h = layer["bbox"]
        accent = rgb(theme, layer["accent"])
        fill = "FFFFFF"
        if layer["accent"] == "green":
            fill = theme["green_light"]
        elif layer["accent"] == "purple":
            fill = theme["purple_light"]
        add_rect(slide, left=px_to_emu_x(x), top=px_to_emu_y(y), width=px_to_emu_x(w), height=px_to_emu_y(h), fill_hex=fill, line_hex=accent, line_pt=0.8, radius=0.08)
        line(slide, x + 205, y, x + 205, y + h, theme["card_border"], 0.8)
        text_box(slide, [x + 22, y + 27, 190, 22], layer["name"], size=9, color=accent, font=fonts["body"], bold=True)
        text_box(slide, [x + 22, y + 57, 180, 44], layer["body"], size=7, color=theme["body_color"], font=fonts["body"])
        if layer.get("middle_label"):
            text_box(slide, [842, y + 52, 175, 22], layer["middle_label"], size=7, color=theme["body_color"], font=fonts["body"], align=PP_ALIGN.CENTER)
        if layer.get("right_label"):
            text_box(slide, [1117, y + 25, 165, 22], layer["right_label"], size=7, color=theme["body_color"], font=fonts["body"], align=PP_ALIGN.CENTER)
        logos = layer.get("logos", [])
        logo_area_x = x + 230
        logo_area_w = max(120, w - 250)
        slot_w = max(70, logo_area_w / max(1, len(logos)))
        for idx, logo in enumerate(logos):
            dynamic_fallback = [
                int(logo_area_x + idx * slot_w + 10),
                int(y + h / 2 - 18),
                int(max(50, slot_w - 20)),
                36,
            ]
            fallback = asset_or_fallback_bbox(spec, logo, dynamic_fallback, layer["bbox"])
            render_logo(slide, spec, base_dir, logo, fallback, labels.get(logo, logo), theme, fonts)
        if layer.get("left_caption"):
            line(slide, 814, y + h - 10, 1042, y + h - 10, accent, 0.8)
            text_box(slide, [860, y + h + 13, 140, 22], layer["left_caption"], size=7, color=accent, font=fonts["body"], align=PP_ALIGN.CENTER)
        if layer.get("right_caption"):
            line(slide, 1080, y + h - 10, 1328, y + h - 10, accent, 0.8, dashed=True)
            text_box(slide, [1128, y + h + 13, 160, 22], layer["right_caption"], size=7, color=accent, font=fonts["body"], align=PP_ALIGN.CENTER)

    line(slide, 928, 314, 928, 337, theme["accent_color"], 1.0, dashed=True)
    line(slide, 928, 455, 928, 477, theme["accent_color"], 1.0, dashed=True)


def render_saas(slide, spec, theme, fonts, base_dir):
    panel = spec["saas_panel"]
    x, y, w, h = panel["bbox"]
    add_rect(slide, left=px_to_emu_x(x), top=px_to_emu_y(y), width=px_to_emu_x(w), height=px_to_emu_y(h), fill_hex="FFFFFF", line_hex=theme["card_border"], line_pt=1.0, radius=0.12)
    text_box(slide, [x + 24, y + 17, w - 48, 24], panel["title"], size=9, color=theme["accent_color"], font=fonts["body"], bold=True, align=PP_ALIGN.CENTER)
    labels = {item["name"]: item["match"] for item in spec.get("logo_assets", [])}
    logos = panel.get("logos", [])
    for idx, logo in enumerate(logos):
        dynamic_y = y + 70 + idx * 50
        fallback = asset_or_fallback_bbox(
            spec,
            logo,
            [x + 28, dynamic_y, w - 56, 36],
            panel["bbox"],
        )
        render_logo(slide, spec, base_dir, logo, fallback, labels.get(logo, logo), theme, fonts)
    text_box(slide, [x + 35, y + h - 53, w - 70, 24], panel["footer"], size=8, color=theme["body_color"], font=fonts["body"], align=PP_ALIGN.CENTER)

    for ly in [259, 386, 535]:
        line(slide, 1351, ly, 1452, ly, theme["accent_color"], 1.0, dashed=True)
    line(slide, 1400, 259, 1400, 646, theme["accent_color"], 1.0, dashed=True)


def render_parallel_layer(slide, spec, theme, fonts, base_dir):
    item = spec["parallel_layer"]
    x, y, w, h = item["bbox"]
    add_rect(slide, left=px_to_emu_x(x), top=px_to_emu_y(y), width=px_to_emu_x(w), height=px_to_emu_y(h), fill_hex=theme["accent_color"], line_hex=None, radius=0.08)
    text_box(slide, [x + 22, y + 27, 330, 24], item["title"], size=10, color="FFFFFF", font=fonts["body"], bold=True)
    text_box(slide, [x + 22, y + 61, 275, 42], item["body"], size=7, color="FFFFFF", font=fonts["body"])
    logo_path = base_dir / spec["header"].get("logo_image", "")
    if logo_path.exists() and logo_path.is_file():
        add_contained_picture(slide, logo_path, left=px_to_emu_x(1440), top=px_to_emu_y(731), width=px_to_emu_x(168), height=px_to_emu_y(38))
    cap_xs = [880, 1070, 1265]
    for cx, cap in zip(cap_xs, item["capabilities"]):
        add_rect(slide, left=px_to_emu_x(cx), top=px_to_emu_y(728), width=px_to_emu_x(145), height=px_to_emu_y(60), fill_hex="EAF2FF", line_hex="BBD1FF", line_pt=1.0, radius=0.2)
        circle_badge(slide, cx + 28, 758, 42, theme, "target" if cap["icon"] != "browser" else "globe")
        text_box(slide, [cx + 58, 744, 75, 36], cap["title"], size=7, color=theme["accent_color"], font=fonts["body"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for tx in [905, 1035, 1185, 1335, 1522]:
        line(slide, tx, 681, tx, 650, theme["accent_color"], 1.0, dashed=True)
    line(slide, 905, 664, 1400, 664, theme["accent_color"], 1.0, dashed=True)
    line(slide, 1522, 679, 1522, 652, theme["accent_color"], 1.0, dashed=True)


def render_callout(slide, spec, theme, fonts):
    c = spec["callout"]
    x, y, w, h = c["bbox"]
    add_rect(slide, left=px_to_emu_x(x), top=px_to_emu_y(y), width=px_to_emu_x(w), height=px_to_emu_y(h), fill_hex="FFFFFF", line_hex=theme["card_border"], line_pt=1.0, radius=0.12)
    star = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, px_to_emu_x(x + 20), px_to_emu_y(y + 22), px_to_emu_x(42), px_to_emu_y(42))
    set_solid_fill(star, theme["accent_color"])
    no_line(star)
    text_box(slide, [x + 74, y + 20, w - 100, 25], c["title"], size=11, color=theme["title_color"], font=fonts["body"], bold=True)
    text_box(slide, [x + 74, y + 49, w - 100, 25], c["body"], size=9, color=theme["body_color"], font=fonts["body"])


def render_slide(spec: dict, base_dir: Path) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = spec["theme"]
    fonts = spec.get("fonts", {"title": "Helvetica Neue", "body": "Helvetica Neue"})

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_solid_fill(bg, theme["bg_color"])
    no_line(bg)

    render_header(slide, spec, theme, fonts, base_dir)
    render_left_panel(slide, spec, theme, fonts)
    render_devices(slide, spec, theme, fonts)
    render_layers(slide, spec, theme, fonts, base_dir)
    render_saas(slide, spec, theme, fonts, base_dir)
    render_parallel_layer(slide, spec, theme, fonts, base_dir)
    render_callout(slide, spec, theme, fonts)
    return prs


def main() -> None:
    if len(sys.argv) < 3:
        print("usage: render_architecture.py <spec.json> <out.pptx>", file=sys.stderr)
        sys.exit(2)
    spec_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2])
    spec = json.loads(spec_path.read_text())
    prs = render_slide(spec, spec_path.parent.parent)
    prs.save(out_path)
    print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
