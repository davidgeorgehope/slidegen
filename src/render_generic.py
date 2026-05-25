"""Render generic slide specs as editable PowerPoint slides."""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

from pptx_utils import (
    DEFAULT_SOURCE_SIZE,
    SLIDE_H,
    SLIDE_W,
    add_contained_picture,
    add_rect,
    add_text,
    hex_to_rgb,
    no_line,
    px_box,
    px_to_emu_x,
    px_to_emu_y,
    set_solid_fill,
    spec_asset_path,
)


ROOT = Path(__file__).resolve().parents[1]
HEX_RE = re.compile(r"^[0-9a-fA-F]{6}$")
DEFAULT_THEME = {
    "bg_color": "F7FBFE",
    "panel_fill": "FFFFFF",
    "title_color": "051353",
    "accent_color": "285CDD",
    "body_color": "051353",
    "muted_color": "5A6A8A",
    "card_border": "D5E3EF",
    "green": "1C9B8A",
    "green_light": "F0FBF8",
    "purple": "7047D6",
    "purple_light": "F7F3FF",
    "blue_light": "EAF2FF",
}
DEFAULT_FONTS = {"title": "Helvetica Neue", "body": "Helvetica Neue"}


def source_size(spec: dict[str, Any]) -> tuple[int, int]:
    value = spec.get("source_size") or DEFAULT_SOURCE_SIZE
    if isinstance(value, list) and len(value) == 2:
        try:
            return int(value[0]), int(value[1])
        except (TypeError, ValueError):
            pass
    return DEFAULT_SOURCE_SIZE


def theme_value(theme: dict[str, str], value: Any, fallback: str | None = None) -> str | None:
    if fallback and fallback in theme:
        fallback = theme[fallback]
    if value is None:
        return fallback
    text = str(value).strip()
    if not text or text.lower() in {"none", "transparent", "null"}:
        return None
    if text in theme:
        text = theme[text]
    text = text.lstrip("#")
    if HEX_RE.match(text):
        return text.upper()
    return fallback


def bool_value(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "y"}
    return bool(value)


def align_value(value: Any):
    text = str(value or "left").lower()
    if text == "center":
        return PP_ALIGN.CENTER
    if text == "right":
        return PP_ALIGN.RIGHT
    if text == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def anchor_value(value: Any):
    text = str(value or "top").lower()
    if text in {"middle", "center"}:
        return MSO_ANCHOR.MIDDLE
    if text == "bottom":
        return MSO_ANCHOR.BOTTOM
    return MSO_ANCHOR.TOP


def role_font_size(role: str, bbox: list[int]) -> int:
    _, _, _, h = bbox
    role = role.lower()
    if role in {"title", "headline", "h1"}:
        return max(22, min(42, int(h * 0.42)))
    if role in {"subtitle", "section", "h2"}:
        return max(15, min(26, int(h * 0.34)))
    if role in {"caption", "eyebrow", "footer", "label"}:
        return max(7, min(12, int(h * 0.34)))
    return max(8, min(18, int(h * 0.32)))


def source_px_to_points(px: float, source: tuple[int, int]) -> float:
    return px * (7.5 * 72) / source[1]


def visual_text_units(text: str) -> float:
    units = 0.0
    for char in text:
        if char.isspace():
            units += 0.35
        elif char in ".,;:!/()[]{}|":
            units += 0.45
        elif char in "-–—":
            units += 0.55
        else:
            units += 1.0
    return max(1.0, units)


def raw_font_size(element: dict[str, Any], source: tuple[int, int]) -> float:
    if element.get("font_size") is not None:
        try:
            return source_px_to_points(float(element["font_size"]), source)
        except (TypeError, ValueError):
            pass
    return source_px_to_points(role_font_size(str(element.get("role") or "body"), element["bbox"]), source)


def fitted_font_size(element: dict[str, Any], source: tuple[int, int]) -> int:
    bbox = element["bbox"]
    text = str(element.get("text") or "")
    explicit_lines = max(1, text.count("\n") + 1)
    raw_pt = raw_font_size(element, source)

    box_height_pt = source_px_to_points(float(bbox[3]), source)
    role = str(element.get("role") or "body").lower()
    box_width_pt = source_px_to_points(float(bbox[2]), source)
    avg_char_width = 0.52 if role in {"title", "headline", "h1", "subtitle", "section", "h2", "heading"} else 0.46
    width_caps = []
    for line in text.splitlines() or [text]:
        width_caps.append(box_width_pt / max(1.0, visual_text_units(line) * avg_char_width))
    max_width_pt = min(width_caps) if width_caps else raw_pt
    max_height_pt = box_height_pt / (explicit_lines * 1.15)
    if role in {"title", "headline", "h1"}:
        max_role_pt = 44
        min_pt = 18
    elif role in {"subtitle", "section", "h2", "heading"}:
        max_role_pt = 24
        min_pt = 10
    elif role in {"caption", "eyebrow", "footer", "label"}:
        max_role_pt = 12
        min_pt = 6
    else:
        max_role_pt = 16
        min_pt = 7
    size = min(raw_pt, max_width_pt, max_height_pt, max_role_pt)
    return int(round(max(min_pt, size)))


def grouped_text_elements(elements: list[dict[str, Any]], source: tuple[int, int]) -> list[dict[str, Any]]:
    """Keep related display-text fragments at one size.

    GPT often emits a colored headline as separate text boxes so individual
    words can have different colors. Fitting each box independently makes the
    line visually inconsistent, so related display text shares the smallest
    fitted size in the group.
    """
    result = []
    groups: dict[tuple[Any, ...], list[dict[str, Any]]] = {}
    groupable_roles = {
        "title",
        "headline",
        "h1",
        "subtitle",
        "section",
        "h2",
        "heading",
        "section_heading",
        "section_title",
        "kicker",
    }
    for element in elements:
        item = dict(element)
        if item.get("type") != "text" or not item.get("bbox"):
            result.append(item)
            continue
        role = str(item.get("role") or "body").lower()
        item["_render_font_size"] = fitted_font_size(item, source)
        if role in groupable_roles and item.get("font_size") is not None:
            key = (
                role,
                str(item.get("font_size")),
                bool_value(item.get("bold")),
                bool_value(item.get("italic")),
                int(item["bbox"][0] // 420),
            )
            groups.setdefault(key, []).append(item)
        result.append(item)

    for items in groups.values():
        if len(items) < 2:
            continue
        size = min(item["_render_font_size"] for item in items)
        raw_sizes = [raw_font_size(item, source) for item in items]
        # If every fragment fits near its raw size, snap the whole group back to
        # the common raw size instead of preserving tiny one-point differences.
        raw_min = min(raw_sizes)
        if all(item["_render_font_size"] >= raw_min - 1 for item in items):
            size = int(round(raw_min))
        for item in items:
            item["_render_font_size"] = size
    return result


def line_shape(slide, points: list[int], theme: dict[str, str], source: tuple[int, int], element: dict[str, Any]):
    x1, y1, x2, y2 = points
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        px_to_emu_x(x1, source),
        px_to_emu_y(y1, source),
        px_to_emu_x(x2, source),
        px_to_emu_y(y2, source),
    )
    color = theme_value(theme, element.get("stroke"), "accent_color") or theme["accent_color"]
    shape.line.color.rgb = hex_to_rgb(color)
    shape.line.width = Pt(float(element.get("stroke_width", 1.0) or 1.0))
    if bool_value(element.get("dash")):
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if bool_value(element.get("arrow")):
        try:
            shape.line.end_arrowhead = True
        except Exception:
            pass
    return shape


def render_text(slide, element: dict[str, Any], theme: dict[str, str], fonts: dict[str, str], source: tuple[int, int]):
    bbox = element["bbox"]
    role = str(element.get("role") or "body")
    color = theme_value(theme, element.get("color"), "body_color") or theme["body_color"]
    size = int(element.get("_render_font_size") or fitted_font_size(element, source))
    font = fonts.get("title" if role.lower() in {"title", "headline", "h1", "section", "h2"} else "body", "Helvetica Neue")
    left, top, width, height = px_box(bbox, source)
    return add_text(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        text=str(element.get("text") or ""),
        font_size=size,
        color_hex=color,
        bold=bool_value(element.get("bold")),
        italic=bool_value(element.get("italic")),
        align=align_value(element.get("align")),
        anchor=anchor_value(element.get("valign")),
        font_name=font,
    )


def render_shape(slide, element: dict[str, Any], theme: dict[str, str], source: tuple[int, int]):
    left, top, width, height = px_box(element["bbox"], source)
    fill = theme_value(theme, element.get("fill"), None)
    stroke = theme_value(theme, element.get("stroke"), None)
    stroke_width = float(element.get("stroke_width", 0.75) or 0.75)
    shape_name = str(element.get("shape") or "rect").lower()
    shape_type = {
        "ellipse": MSO_SHAPE.OVAL,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    }.get(shape_name, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        set_solid_fill(shape, fill)
    if stroke is None or stroke_width <= 0:
        no_line(shape)
    else:
        shape.line.color.rgb = hex_to_rgb(stroke)
        shape.line.width = Pt(stroke_width)
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shape.adjustments[0] = float(element.get("radius", 0.12) or 0.12)
        except Exception:
            pass
    if element.get("text"):
        shape.text_frame.text = ""
    return shape


def render_image(slide, element: dict[str, Any], spec: dict[str, Any], base_dir: Path, theme: dict[str, str], source: tuple[int, int]):
    path = spec_asset_path(spec, base_dir, element.get("asset"))
    left, top, width, height = px_box(element["bbox"], source)
    if path:
        return add_contained_picture(slide, path, left=left, top=top, width=width, height=height)

    if not bool_value(element.get("placeholder")):
        return None

    add_rect(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        fill_hex=theme_value(theme, element.get("fill"), "panel_fill"),
        line_hex=theme_value(theme, element.get("stroke"), "card_border"),
        line_pt=0.75,
        radius=0.08,
    )
    return None


def icon_line(slide, x1, y1, x2, y2, color, width, source):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        px_to_emu_x(x1, source),
        px_to_emu_y(y1, source),
        px_to_emu_x(x2, source),
        px_to_emu_y(y2, source),
    )
    shape.line.color.rgb = hex_to_rgb(color)
    shape.line.width = Pt(width)
    return shape


def icon_shape(slide, shape_type, x, y, w, h, fill, stroke, source, line_width=1.2):
    shape = slide.shapes.add_shape(
        shape_type,
        px_to_emu_x(x, source),
        px_to_emu_y(y, source),
        px_to_emu_x(w, source),
        px_to_emu_y(h, source),
    )
    if fill:
        set_solid_fill(shape, fill)
    else:
        shape.fill.background()
    if stroke:
        shape.line.color.rgb = hex_to_rgb(stroke)
        shape.line.width = Pt(line_width)
    else:
        no_line(shape)
    return shape


def draw_symbol(slide, bbox: list[int], hint: str, accent: str, source: tuple[int, int]):
    x, y, w, h = bbox
    cx = x + w / 2
    cy = y + h / 2
    size = min(w, h)
    hint_l = hint.lower()
    white = "FFFFFF"
    stroke = white
    sw = max(1.0, size / 32)

    if any(word in hint_l for word in ("lock", "secure", "security", "privacy")):
        icon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx - size * 0.20, cy - size * 0.02, size * 0.40, size * 0.28, white, None, source)
        arc = slide.shapes.add_shape(
            MSO_SHAPE.ARC,
            px_to_emu_x(cx - size * 0.18, source),
            px_to_emu_y(cy - size * 0.30, source),
            px_to_emu_x(size * 0.36, source),
            px_to_emu_y(size * 0.38, source),
        )
        arc.line.color.rgb = hex_to_rgb(stroke)
        arc.line.width = Pt(sw * 1.5)
        arc.fill.background()
        return

    if any(word in hint_l for word in ("shield", "protect", "govern", "policy", "compliance")):
        icon_shape(slide, MSO_SHAPE.PENTAGON, cx - size * 0.22, cy - size * 0.28, size * 0.44, size * 0.52, white, None, source)
        return

    if any(word in hint_l for word in ("user", "identity", "people", "group", "role")):
        offsets = (-size * 0.13, size * 0.13) if any(word in hint_l for word in ("users", "people", "group")) else (0,)
        for dx in offsets:
            icon_shape(slide, MSO_SHAPE.OVAL, cx + dx - size * 0.08, cy - size * 0.26, size * 0.16, size * 0.16, white, None, source)
            icon_shape(slide, MSO_SHAPE.OVAL, cx + dx - size * 0.13, cy - size * 0.06, size * 0.26, size * 0.25, None, white, source, sw)
        return

    if any(word in hint_l for word in ("database", "data", "storage", "warehouse")):
        shape_type = getattr(MSO_SHAPE, "CAN", MSO_SHAPE.OVAL)
        icon_shape(slide, shape_type, cx - size * 0.22, cy - size * 0.26, size * 0.44, size * 0.52, None, white, source, sw)
        icon_line(slide, cx - size * 0.18, cy - size * 0.05, cx + size * 0.18, cy - size * 0.05, white, sw, source)
        icon_line(slide, cx - size * 0.18, cy + size * 0.12, cx + size * 0.18, cy + size * 0.12, white, sw, source)
        return

    if any(word in hint_l for word in ("cloud", "saas", "app", "application")):
        shape_type = getattr(MSO_SHAPE, "CLOUD", MSO_SHAPE.OVAL)
        icon_shape(slide, shape_type, cx - size * 0.28, cy - size * 0.18, size * 0.56, size * 0.38, None, white, source, sw)
        return

    if any(word in hint_l for word in ("browser", "window", "web", "portal")):
        icon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx - size * 0.28, cy - size * 0.23, size * 0.56, size * 0.44, None, white, source, sw)
        icon_line(slide, cx - size * 0.28, cy - size * 0.10, cx + size * 0.28, cy - size * 0.10, white, sw, source)
        return

    if any(word in hint_l for word in ("laptop", "device", "endpoint", "computer")):
        icon_shape(slide, MSO_SHAPE.RECTANGLE, cx - size * 0.28, cy - size * 0.22, size * 0.56, size * 0.34, None, white, source, sw)
        icon_line(slide, cx - size * 0.34, cy + size * 0.18, cx + size * 0.34, cy + size * 0.18, white, sw, source)
        return

    if "phone" in hint_l or "mobile" in hint_l:
        icon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx - size * 0.14, cy - size * 0.30, size * 0.28, size * 0.60, None, white, source, sw)
        return

    if any(word in hint_l for word in ("alert", "risk", "warning", "threat")):
        icon_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, cx - size * 0.25, cy - size * 0.24, size * 0.50, size * 0.50, None, white, source, sw)
        icon_line(slide, cx, cy - size * 0.10, cx, cy + size * 0.08, white, sw, source)
        icon_shape(slide, MSO_SHAPE.OVAL, cx - size * 0.018, cy + size * 0.16, size * 0.036, size * 0.036, white, None, source)
        return

    if any(word in hint_l for word in ("sync", "flow", "connector", "integration", "workflow", "pipeline")):
        icon_line(slide, cx - size * 0.28, cy - size * 0.06, cx + size * 0.20, cy - size * 0.06, white, sw, source)
        icon_line(slide, cx + size * 0.14, cy - size * 0.14, cx + size * 0.24, cy - size * 0.06, white, sw, source)
        icon_line(slide, cx + size * 0.14, cy + size * 0.02, cx + size * 0.24, cy - size * 0.06, white, sw, source)
        icon_line(slide, cx + size * 0.28, cy + size * 0.08, cx - size * 0.20, cy + size * 0.08, white, sw, source)
        icon_line(slide, cx - size * 0.14, cy, cx - size * 0.24, cy + size * 0.08, white, sw, source)
        icon_line(slide, cx - size * 0.14, cy + size * 0.16, cx - size * 0.24, cy + size * 0.08, white, sw, source)
        return

    # Unknown semantic icon: draw an abstract target/radar mark instead of text.
    for scale in (0.55, 0.32, 0.10):
        d = size * scale
        icon_shape(slide, MSO_SHAPE.OVAL, cx - d / 2, cy - d / 2, d, d, None if scale != 0.10 else white, white, source, sw)


def render_icon(slide, element: dict[str, Any], spec: dict[str, Any], base_dir: Path, theme: dict[str, str], source: tuple[int, int]):
    path = spec_asset_path(spec, base_dir, element.get("asset"))
    left, top, width, height = px_box(element["bbox"], source)
    if path:
        return add_contained_picture(slide, path, left=left, top=top, width=width, height=height)

    x, y, w, h = element["bbox"]
    accent = theme_value(theme, element.get("color"), "accent_color") or theme["accent_color"]
    outer = add_rect(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        fill_hex=theme_value(theme, element.get("background"), "blue_light"),
        line_hex=theme_value(theme, element.get("stroke"), "card_border"),
        line_pt=0.75,
        radius=0.5,
    )
    try:
        outer.adjustments[0] = 0.5
    except Exception:
        pass
    inset = max(4, int(min(w, h) * 0.17))
    icon_shape(
        slide,
        MSO_SHAPE.OVAL,
        x + inset,
        y + inset,
        max(1, w - inset * 2),
        max(1, h - inset * 2),
        accent,
        None,
        source,
    )
    draw_symbol(slide, [x + inset, y + inset, max(1, w - inset * 2), max(1, h - inset * 2)], str(element.get("icon_hint") or element.get("name") or ""), accent, source)
    return outer


def add_slide(prs: Presentation, spec: dict[str, Any], base_dir: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = dict(DEFAULT_THEME)
    if isinstance(spec.get("theme"), dict):
        theme.update({str(k): str(v) for k, v in spec["theme"].items() if v is not None})
    fonts = dict(DEFAULT_FONTS)
    if isinstance(spec.get("fonts"), dict):
        fonts.update({str(k): str(v) for k, v in spec["fonts"].items() if v})
    source = source_size(spec)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_solid_fill(bg, theme_value(theme, theme.get("bg_color"), "bg_color") or theme["bg_color"])
    no_line(bg)

    for element in grouped_text_elements(spec.get("elements", []), source):
        if not isinstance(element, dict):
            continue
        element_type = element.get("type")
        if element_type == "shape" and element.get("bbox"):
            render_shape(slide, element, theme, source)
        elif element_type == "line" and element.get("points"):
            line_shape(slide, element["points"], theme, source, element)
        elif element_type == "image" and element.get("bbox"):
            render_image(slide, element, spec, base_dir, theme, source)
        elif element_type == "icon" and element.get("bbox"):
            render_icon(slide, element, spec, base_dir, theme, source)
        elif element_type == "text" and element.get("bbox"):
            render_text(slide, element, theme, fonts, source)
    return slide


def specs_from_root(spec: dict[str, Any]) -> list[dict[str, Any]]:
    if spec.get("layout") != "generic_deck":
        return [spec]

    shared = {
        "source_size": spec.get("source_size"),
        "theme": spec.get("theme", {}),
        "fonts": spec.get("fonts", {}),
        "header": spec.get("header", {}),
        "logo_assets": spec.get("logo_assets", []),
        "asset_queries": spec.get("asset_queries", []),
        "assets": spec.get("assets", {}),
    }
    slides = []
    for idx, child in enumerate(spec.get("slides", []), start=1):
        if not isinstance(child, dict):
            continue
        merged = dict(shared)
        merged.update(child)
        merged["layout"] = "generic_slide"
        merged["source_size"] = child.get("source_size") or shared["source_size"]
        merged["theme"] = {**shared.get("theme", {}), **child.get("theme", {})}
        merged["fonts"] = {**shared.get("fonts", {}), **child.get("fonts", {})}
        merged["header"] = {**shared.get("header", {}), **child.get("header", {})}
        merged["assets"] = {**shared.get("assets", {}), **child.get("assets", {})}
        merged["slide"] = child.get("slide") or f"{spec.get('slide', 'slide')}_{idx:02d}"
        slides.append(merged)
    return slides


def render_deck(specs: list[dict[str, Any]], base_dir: Path) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for spec in specs:
        add_slide(prs, spec, base_dir)
    return prs


def render_slide(spec: dict[str, Any], base_dir: Path) -> Presentation:
    return render_deck(specs_from_root(spec), base_dir)


def main() -> None:
    if len(sys.argv) < 3:
        print("usage: render_generic.py <spec.json> <out.pptx>", file=sys.stderr)
        sys.exit(2)
    spec_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2])
    spec = json.loads(spec_path.read_text())
    prs = render_slide(spec, ROOT)
    prs.save(out_path)
    print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
