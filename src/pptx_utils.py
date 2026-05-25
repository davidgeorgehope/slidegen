"""Shared PowerPoint rendering helpers."""
from __future__ import annotations

from pathlib import Path

from PIL import Image as PILImage
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
DEFAULT_SOURCE_SIZE = (1672, 941)


def hex_to_rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(str(hex_str).lstrip("#"))


def set_solid_fill(shape, hex_color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(hex_color)


def no_line(shape) -> None:
    shape.line.fill.background()


def emu(value) -> Emu:
    return Emu(int(value))


def px_to_emu_x(px: float, source_size=DEFAULT_SOURCE_SIZE):
    return emu(px / source_size[0] * SLIDE_W)


def px_to_emu_y(px: float, source_size=DEFAULT_SOURCE_SIZE):
    return emu(px / source_size[1] * SLIDE_H)


def px_box(bbox, source_size=DEFAULT_SOURCE_SIZE):
    x, y, w, h = bbox
    return (
        px_to_emu_x(x, source_size),
        px_to_emu_y(y, source_size),
        px_to_emu_x(w, source_size),
        px_to_emu_y(h, source_size),
    )


def add_rect(
    slide,
    *,
    left,
    top,
    width,
    height,
    fill_hex=None,
    line_hex=None,
    line_pt=0.75,
    radius=None,
    no_fill=False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if no_fill:
        shape.fill.background()
    elif fill_hex is not None:
        set_solid_fill(shape, fill_hex)
    if line_hex is None:
        no_line(shape)
    else:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(line_pt)
    if radius is not None:
        shape.adjustments[0] = radius
    shape.text_frame.text = ""
    return shape


def add_contained_picture(slide, img_path: Path, *, left, top, width, height):
    """Add an image scaled to fit inside a box while preserving aspect ratio."""
    with PILImage.open(img_path) as pim:
        iw, ih = pim.size

    box_aspect = width / height
    img_aspect = iw / ih
    if img_aspect > box_aspect:
        new_w = width
        new_h = emu(width / img_aspect)
    else:
        new_h = height
        new_w = emu(height * img_aspect)

    pic_x = emu(left + (width - new_w) / 2)
    pic_y = emu(top + (height - new_h) / 2)
    return slide.shapes.add_picture(str(img_path), pic_x, pic_y, width=new_w, height=new_h)


def add_picture_bbox(slide, img_path: Path, bbox, source_size=DEFAULT_SOURCE_SIZE):
    left, top, width, height = px_box(bbox, source_size)
    return slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)


def spec_asset_path(spec: dict, base_dir: Path, name: str | None) -> Path | None:
    if not name:
        return None
    asset = spec.get("assets", {}).get(name)
    if not asset or not asset.get("path"):
        return None
    path = base_dir / asset["path"]
    return path if path.exists() else None


def set_textframe_padding(tf, pad_in=0.0):
    tf.margin_left = tf.margin_right = emu(Inches(pad_in))
    tf.margin_top = tf.margin_bottom = emu(Inches(pad_in))


def add_text(
    slide,
    *,
    left,
    top,
    width,
    height,
    text: str,
    font_size: int,
    color_hex: str,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font_name: str = "Helvetica Neue",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    set_textframe_padding(tf)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = hex_to_rgb(color_hex)
    return tb
