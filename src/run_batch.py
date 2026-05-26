"""Run the slide rebuild pipeline over a directory of source images."""
from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation

from render_generic import SLIDE_H, SLIDE_W, add_slide, specs_from_root


ROOT = Path(__file__).resolve().parents[1]


def load_env(path: Path) -> dict[str, str]:
    env = os.environ.copy()
    if not path.exists():
        return env
    for raw in path.read_text().splitlines():
        line = raw.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        key = key.strip()
        value = value.strip().strip('"').strip("'")
        if key:
            env.setdefault(key, value)
    return env


def natural_key(path: Path):
    parts = re.split(r"(\d+)", path.stem)
    return [int(part) if part.isdigit() else part.lower() for part in parts]


def discover_images(images_dir: Path, pattern: str, start: int | None, end: int | None) -> list[Path]:
    images = sorted(images_dir.glob(pattern), key=natural_key)
    if start is None and end is None:
        return images

    filtered = []
    for path in images:
        match = re.search(r"(\d+)$", path.stem)
        if not match:
            continue
        number = int(match.group(1))
        if start is not None and number < start:
            continue
        if end is not None and number > end:
            continue
        filtered.append(path)
    return filtered


def parse_slide_numbers(value: str | None) -> set[int] | None:
    if not value:
        return None
    numbers = set()
    for raw in value.split(","):
        item = raw.strip()
        if not item:
            continue
        if "-" in item:
            left, right = item.split("-", 1)
            start = int(left)
            end = int(right)
            numbers.update(range(min(start, end), max(start, end) + 1))
        else:
            numbers.add(int(item))
    return numbers or None


def run(cmd: list[str], env: dict[str, str], dry_run: bool) -> None:
    print("+ " + " ".join(cmd), flush=True)
    if dry_run:
        return
    subprocess.run(cmd, cwd=ROOT, env=env, check=True)


def pipeline_cmd(args, image_path: Path, spec_path: Path, out_path: Path) -> list[str]:
    cmd = [
        sys.executable,
        "src/run_pipeline.py",
        str(image_path),
        str(spec_path),
        str(out_path),
        "--asset-mode",
        args.asset_mode,
        "--spec-layout",
        args.spec_layout,
    ]
    if args.spec_model:
        cmd.extend(["--spec-model", args.spec_model])
    for helper in args.style_guide_image:
        cmd.extend(["--style-guide-image", helper])
    if args.template_pptx:
        cmd.extend(["--template-pptx", args.template_pptx])
    if args.extract_root:
        cmd.extend(["--extract-root", args.extract_root])
    if args.icon_library_dir:
        cmd.extend(["--icon-library-dir", args.icon_library_dir])
    if args.refine:
        cmd.extend(["--refine", str(args.refine)])
    if args.refine_model:
        cmd.extend(["--refine-model", args.refine_model])
    if args.force_spec:
        cmd.extend(["--generate-spec", "--force-spec"])
    elif not spec_path.exists():
        cmd.append("--generate-spec")
    if args.skip_assets:
        cmd.append("--skip-assets")
    if args.skip_generic_assets:
        cmd.append("--skip-generic-assets")
    if args.no_verify:
        cmd.append("--no-verify")
    return cmd


def combine_specs(spec_paths: list[Path], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    added = 0
    skipped = []
    for spec_path in spec_paths:
        spec = json.loads(spec_path.read_text())
        if spec.get("layout") not in {"generic_slide", "generic_deck"}:
            skipped.append(spec_path)
            continue
        for slide_spec in specs_from_root(spec):
            add_slide(prs, slide_spec, ROOT)
            added += 1

    if added == 0:
        raise RuntimeError("No generic slides were available for the combined deck")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote {out_path} with {added} slide(s)")
    if skipped:
        print("skipped non-generic specs in combined deck:")
        for path in skipped:
            print(f"- {path}")


def render_preview_cmd(pptx_path: Path, png_path: Path) -> list[str]:
    return [sys.executable, "src/render_preview.py", str(pptx_path), str(png_path)]


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--images-dir", default="images")
    parser.add_argument("--pattern", default="image*.png")
    parser.add_argument("--spec-dir", default="specs/auto")
    parser.add_argument("--output-dir", default="output/auto")
    parser.add_argument("--combined", default="output/all_slides_auto.pptx")
    parser.add_argument("--asset-mode", choices=["auto", "extract", "generate"], default="auto")
    parser.add_argument("--spec-layout", default="generic_slide")
    parser.add_argument("--spec-model", default=None)
    parser.add_argument("--style-guide-image", action="append", default=[])
    parser.add_argument("--template-pptx", default=None)
    parser.add_argument("--refine", type=int, default=2)
    parser.add_argument("--refine-model", default=None)
    parser.add_argument("--start", type=int, default=None)
    parser.add_argument("--end", type=int, default=None)
    parser.add_argument("--slides", default=None, help="comma/range slide numbers, for example 1,26-28")
    parser.add_argument("--limit", type=int, default=None)
    parser.add_argument("--scratch-dir", default=None, help="scratch root for specs/assets/icon cache")
    parser.add_argument("--keep-scratch", action="store_true")
    parser.add_argument("--force-spec", action="store_true")
    parser.add_argument("--skip-assets", action="store_true")
    parser.add_argument(
        "--skip-generic-assets",
        action="store_true",
        help="layout-QA only: extract real logos but draw generic icons as native placeholders",
    )
    parser.add_argument("--no-verify", action="store_true")
    parser.add_argument("--no-combined", action="store_true")
    parser.add_argument("--no-previews", action="store_true", help="skip final PNG previews beside output PPTX files")
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()

    scratch_dir = Path(args.scratch_dir) if args.scratch_dir else None
    if scratch_dir:
        scratch_dir = scratch_dir.resolve()
        if not args.keep_scratch and not args.dry_run:
            shutil.rmtree(scratch_dir, ignore_errors=True)
        args.extract_root = str(scratch_dir / "extracted")
        args.icon_library_dir = str(scratch_dir / "icon_library")
        if args.spec_dir == parser.get_default("spec_dir"):
            args.spec_dir = str(scratch_dir / "specs")
    else:
        args.extract_root = None
        args.icon_library_dir = None

    images_dir = ROOT / args.images_dir
    spec_dir = ROOT / args.spec_dir
    output_dir = ROOT / args.output_dir
    combined_path = ROOT / args.combined
    images = discover_images(images_dir, args.pattern, args.start, args.end)
    slide_numbers = parse_slide_numbers(args.slides)
    if slide_numbers is not None:
        selected = []
        for image in images:
            match = re.search(r"(\d+)$", image.stem)
            if match and int(match.group(1)) in slide_numbers:
                selected.append(image)
        images = selected
    if args.limit is not None:
        images = images[: args.limit]
    if not images:
        raise SystemExit(f"No images matched {images_dir / args.pattern}")

    env = load_env(ROOT / ".env")
    spec_dir.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)

    spec_paths = []
    out_paths = []
    for idx, image_path in enumerate(images, start=1):
        spec_path = spec_dir / f"{image_path.stem}.json"
        out_path = output_dir / f"{image_path.stem}.pptx"
        print(f"\n[{idx}/{len(images)}] {image_path.name}", flush=True)
        run(pipeline_cmd(args, image_path, spec_path, out_path), env, args.dry_run)
        spec_paths.append(spec_path)
        out_paths.append(out_path)
        if not args.no_previews:
            run(render_preview_cmd(out_path, out_path.with_suffix(".png")), env, args.dry_run)

    if not args.no_combined and not args.dry_run:
        combine_specs(spec_paths, combined_path)
        if not args.no_previews:
            run(render_preview_cmd(combined_path, combined_path.with_suffix(".png")), env, args.dry_run)

    if scratch_dir and not args.keep_scratch and not args.dry_run:
        shutil.rmtree(scratch_dir, ignore_errors=True)


if __name__ == "__main__":
    main()
